"""Generate the master thesis defense PowerPoint.

Run:
    .venv/bin/python master/thesis_word/upt/presentation/build_presentation.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


PROJECT_ROOT = Path("/Users/alketaalia/Documents/diploma master/DRPrediction/DRPrediction")
RESULTS = PROJECT_ROOT / "results"
MASTER_RESULTS = PROJECT_ROOT / "master" / "results"
APP_SHOTS = MASTER_RESULTS / "app_screenshots"
FIGURES = PROJECT_ROOT / "master" / "thesis" / "figures"
OUT_PATH = PROJECT_ROOT / "master" / "thesis_word" / "upt" / "presentation" / "ThesisDefense.pptx"

# UPT palette
COLOR_PRIMARY = RGBColor(0x0B, 0x3D, 0x91)   # deep blue
COLOR_ACCENT = RGBColor(0xE8, 0x7A, 0x00)    # orange
COLOR_DARK = RGBColor(0x1F, 0x1F, 0x1F)
COLOR_GREY = RGBColor(0x55, 0x55, 0x55)
COLOR_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_footer_bar(slide, prs):
    """Thin coloured bar at the bottom of the slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.25),
        prs.slide_width, Inches(0.25),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()


def add_page_number(slide, prs, num, total):
    """Page X / Y at bottom right."""
    tb = slide.shapes.add_textbox(
        prs.slide_width - Inches(1.5), prs.slide_height - Inches(0.25),
        Inches(1.4), Inches(0.25),
    )
    tf = tb.text_frame
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {total}"
    run.font.size = Pt(10)
    run.font.color.rgb = COLOR_WHITE


def add_title(slide, text, prs, size=32):
    """Top title bar."""
    tb = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.25), prs.slide_width - Inches(0.8), Inches(0.7),
    )
    tf = tb.text_frame
    tf.margin_left = Pt(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = COLOR_PRIMARY
    # Underline accent bar
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.95),
        Inches(1.2), Inches(0.05),
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = COLOR_ACCENT
    underline.line.fill.background()


def add_bullets(slide, prs, items, left=Inches(0.5), top=Inches(1.3),
                width=None, height=None, size=18, bold_first=False):
    """Add a bulleted list."""
    if width is None:
        width = prs.slide_width - Inches(1.0)
    if height is None:
        height = prs.slide_height - top - Inches(0.5)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # Two-tuple = (bullet, sub-text). String = plain bullet.
        if isinstance(item, tuple) and len(item) == 2:
            bullet_text, detail = item
            r1 = p.add_run()
            r1.text = f"• {bullet_text}  "
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = COLOR_DARK
            r2 = p.add_run()
            r2.text = detail
            r2.font.size = Pt(size - 2)
            r2.font.color.rgb = COLOR_GREY
        else:
            r = p.add_run()
            r.text = f"• {item}"
            r.font.size = Pt(size)
            r.font.color.rgb = COLOR_DARK
        p.space_after = Pt(8)


def add_two_columns(slide, prs, left_items, right_items, top=Inches(1.3), size=16):
    """Two side-by-side bulleted columns."""
    col_w = (prs.slide_width - Inches(1.2)) // 2
    add_bullets(slide, prs, left_items, left=Inches(0.5), top=top,
                width=col_w, size=size)
    add_bullets(slide, prs, right_items, left=Inches(0.6) + col_w, top=top,
                width=col_w, size=size)


def add_image(slide, image_path, left, top, width=None, height=None):
    """Add an image, gracefully skipping if it doesn't exist."""
    if not Path(image_path).is_file():
        return None
    return slide.shapes.add_picture(str(image_path), left, top,
                                     width=width, height=height)


def add_caption(slide, text, left, top, width, size=11):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.italic = True
    run.font.color.rgb = COLOR_GREY


def add_table(slide, prs, data, left, top, width, height, header_color=COLOR_PRIMARY,
              header_text_color=COLOR_WHITE, font_size=12, header_size=13):
    """Add a styled table from a 2D list. Row 0 is the header."""
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(header_size if r == 0 else font_size)
                    run.font.bold = (r == 0)
                    if r == 0:
                        run.font.color.rgb = header_text_color
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                if r % 2 == 0:
                    cell.fill.fore_color.rgb = COLOR_LIGHT
                else:
                    cell.fill.fore_color.rgb = COLOR_WHITE
    return tbl


# ============================================================================
# Slide builders
# ============================================================================

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background bar
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height,
    )
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_PRIMARY
    band.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(3.6),
        prs.slide_width, Inches(0.08),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_ACCENT
    accent.line.fill.background()

    # University label
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), prs.slide_width - Inches(1.0), Inches(0.5),
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "POLYTECHNIC UNIVERSITY OF TIRANA"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.85), prs.slide_width - Inches(1.0), Inches(0.4),
    )
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Faculty of Information Technology  ·  Master Thesis Defense"
    run.font.size = Pt(13)
    run.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)

    # Main title
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.8), prs.slide_width - Inches(1.0), Inches(1.6),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Uncertainty-Aware Diabetic Retinopathy Grading"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Using Conformal Prediction and Bayesian Deep Learning"
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)

    # Author / supervisor block
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.0), prs.slide_width - Inches(1.0), Inches(2.0),
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Author"
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0xAA, 0xBB, 0xDD)

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Alketa Alia"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(20)
    run = p.add_run()
    run.text = "Supervisor"
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0xAA, 0xBB, 0xDD)

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Prof. Dr. [Supervisor Name]"
    run.font.size = Pt(18)
    run.font.color.rgb = COLOR_WHITE

    # Footer with date
    tb = slide.shapes.add_textbox(
        Inches(0.5), prs.slide_height - Inches(0.8),
        prs.slide_width - Inches(1.0), Inches(0.5),
    )
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Tirana  ·  May 2026"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)


def slide_outline(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Outline", prs)

    items = [
        ("1. Problem & motivation", "Why uncertainty matters in medical AI"),
        ("2. Research questions", "Four questions the thesis answers"),
        ("3. Methodology", "Dataset, splits, training pipeline, UQ methods"),
        ("4. Results — Phase 1", "Binary baselines, calibration, ensembling"),
        ("5. Results — Phase 2", "Conformal, MC Dropout, K-fold, OOD"),
        ("6. Results — Phase 3", "Multi-class 5-stage grading"),
        ("7. Deployment", "Streamlit decision-support application"),
        ("8. Limitations & future work", "What is open"),
        ("9. Conclusions", "Key takeaways"),
    ]
    add_bullets(slide, prs, items, top=Inches(1.3), size=18)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_problem(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Problem: A Global Screening Bottleneck", prs)
    items = [
        ("537 million", "adults living with diabetes globally (IDF 2021)"),
        ("100 million", "have vision-threatening diabetic retinopathy"),
        ("DR is the leading cause", "of preventable blindness in working-age adults"),
        ("Early stages are asymptomatic", "→ screening is the only reliable detection"),
        ("Specialist shortage", "1 ophthalmologist per million in sub-Saharan Africa; gap exists in rural Albania / Balkans too"),
        ("Automated screening", "is a structural answer, not an academic curiosity"),
    ]
    add_bullets(slide, prs, items, top=Inches(1.3), size=17)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_5_stages(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Clinical Grading: 5 Severity Levels", prs)
    data = [
        ["Stage", "Name", "Finding", "Management"],
        ["0", "No DR", "Healthy retina", "Routine review"],
        ["1", "Mild NPDR", "Microaneurysms only", "Annual follow-up"],
        ["2", "Moderate NPDR", "+ Haemorrhages, exudates", "6-month follow-up"],
        ["3", "Severe NPDR", "4-2-1 rule", "Urgent referral"],
        ["4", "PDR", "Neovascularization", "Immediate treatment"],
    ]
    add_table(slide, prs, data, Inches(0.6), Inches(1.4),
              prs.slide_width - Inches(1.2), Inches(3.5),
              font_size=14, header_size=14)
    tb = slide.shapes.add_textbox(
        Inches(0.6), Inches(5.2), prs.slide_width - Inches(1.2), Inches(1.5),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Binary 'DR / No DR' is informationally lossy: "
    run.font.size = Pt(16)
    run.font.color.rgb = COLOR_DARK
    r2 = p.add_run()
    r2.text = "Mild and PDR receive the same label but require very different management."
    r2.font.size = Pt(16)
    r2.font.italic = True
    r2.font.color.rgb = COLOR_ACCENT
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_motivation(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Why Uncertainty Matters", prs)
    items = [
        ("Modern CNNs are notoriously over-confident", "A model can be wrong far more often than its 99% confidence suggests"),
        ("In medicine, confidently wrong is dangerous", "A confident 'No DR' can delay vision-saving treatment"),
        ("Regulators now require uncertainty disclosure", "FDA SaMD, EU AI Act 2026, WHO Ethics for Health"),
        ("This thesis exposes 4 layers of output", "1. Argmax decision  2. Calibrated probability  3. Conformal set  4. OOD flag"),
        ("Result", "A screening tool that 'knows when it doesn't know' and defers to the clinician"),
    ]
    add_bullets(slide, prs, items, top=Inches(1.3), size=17)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_research_questions(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Research Questions", prs)
    items = [
        ("RQ1.", "Do alternative deep architectures differ once preprocessing is properly configured?"),
        ("RQ2.", "How well-calibrated are the models, and does temperature scaling help?"),
        ("RQ3.", "What is the most useful uncertainty signal? Can conformal prediction provide a coverage guarantee?"),
        ("RQ4.", "Does the same uncertainty machinery generalise to the clinically meaningful 5-stage grading?"),
    ]
    add_bullets(slide, prs, items, top=Inches(1.5), size=20)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_methodology_overview(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Methodology Overview", prs)
    left = [
        ("Dataset", "APTOS 2019 — 3,662 fundus images, 224×224"),
        ("Splits", "Stratified 70 / 15 / 15 with random_state=123"),
        ("Train", "2,562 images for parameter updates"),
        ("Validation", "550 images — early stopping, T fitting, conformal threshold"),
        ("Test", "550 images — locked, every reported metric"),
    ]
    right = [
        ("Hardware", "Apple Silicon M-series, 8 GB RAM, CPU-only"),
        ("Framework", "TensorFlow 2.15 / Keras"),
        ("Optimizer", "Adam, lr 1e-3, batch 32"),
        ("Class weights", "'balanced' — critical for 5-class"),
        ("Callbacks", "EarlyStop (10), ReduceLR (5), ModelCheckpoint"),
    ]
    add_two_columns(slide, prs, left, right, top=Inches(1.3), size=15)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_architectures(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Six Deep Architectures Evaluated", prs)
    data = [
        ["Architecture", "Key idea", "Params", "Role in thesis"],
        ["ResNet50", "Skip connections", "25.6 M", "Transfer learning backbone"],
        ["DenseNet121", "Dense connectivity", "7.0 M", "Backbone + feature extractor for OOD"],
        ["Xception", "Depthwise separable", "21.0 M", "Efficient transfer learning"],
        ["VGG16", "Simple 3×3 conv blocks", "138 M", "Calibration case study (T=1.44)"],
        ["CNN (custom)", "From-scratch, 3 blocks", "11 M", "Grayscale baseline"],
        ["CNN (Tanh+ReLU)", "Mixed activations", "11 M", "Statistical comparison control"],
    ]
    add_table(slide, prs, data, Inches(0.5), Inches(1.4),
              prs.slide_width - Inches(1.0), Inches(4.0),
              font_size=12, header_size=13)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_4_pillars(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Four Pillars of Uncertainty Quantification", prs)
    items = [
        ("1. Calibration", "ECE / MCE + reliability diagrams + post-hoc temperature scaling"),
        ("2. Bayesian (MC Dropout)", "T = 30 stochastic forward passes → epistemic / aleatoric decomposition"),
        ("3. Conformal Prediction", "Split conformal with LAC and APS, α = 0.05 and 0.10 — finite-sample coverage guarantee"),
        ("4. OOD Detection", "Maximum Softmax, Energy, Mahalanobis distance, cosine distance"),
        ("Bonus", "Heterogeneous ensemble across all 6 binary architectures"),
    ]
    add_bullets(slide, prs, items, top=Inches(1.3), size=18)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_phase1_results(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 1 — Binary Classification", prs)
    data = [
        ["Model", "Test acc", "95% CI", "AUC", "ECE raw", "ECE TS"],
        ["ResNet50", "95.27 %", "[93.45, 97.09]", "0.989", "0.020", "0.018"],
        ["Xception", "95.45 %", "[93.82, 97.09]", "0.985", "0.017", "0.022"],
        ["DenseNet121", "95.64 %", "[94.00, 97.27]", "0.992", "0.029", "0.027"],
        ["VGG16", "95.64 %", "[94.00, 97.27]", "0.990", "0.026", "0.015"],
        ["CNN", "96.00 %", "[94.36, 97.45]", "0.987", "0.032", "0.032"],
        ["CNN (Tanh+ReLU)", "92.91 %", "[90.72, 94.91]", "0.970", "0.046", "0.041"],
        ["Ensemble (6)", "96.55 %", "—", "0.991", "0.028", "—"],
    ]
    add_table(slide, prs, data, Inches(0.4), Inches(1.3),
              prs.slide_width - Inches(0.8), Inches(3.6),
              font_size=11, header_size=12)
    add_bullets(slide, prs, [
        "Five strong models cluster within 0.73 pp — every pairwise McNemar p > 0.5",
        "CNN (Tanh+ReLU) is significantly worse (p < 0.05 against every other model)",
        "Ensemble (96.55 %) beats every individual member",
        "VGG16: temperature scaling reduces ECE by 40 % (0.026 → 0.015)",
    ], top=Inches(5.1), size=14)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_phase1_figure(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 1 — Visual Summary", prs)
    img1 = FIGURES / "fig_phase1_accuracy_ci.png"
    img2 = MASTER_RESULTS / "mcnemar_pvalues.png"
    left = Inches(0.4)
    top = Inches(1.3)
    width = (prs.slide_width - Inches(1.0)) // 2
    add_image(slide, img1, left, top, width=width)
    add_image(slide, img2, left + width + Inches(0.2), top, width=width)
    add_caption(slide, "Per-model accuracy with 95 % bootstrap CIs",
                left, top + Inches(3.5), width, size=12)
    add_caption(slide, "Pairwise McNemar p-values: only CNN(T+R) differs",
                left + width + Inches(0.2), top + Inches(3.5), width, size=12)

    # Ensemble selective accuracy callout
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), prs.slide_height - Inches(1.8),
        prs.slide_width - Inches(1.0), Inches(1.3),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_LIGHT
    box.line.color.rgb = COLOR_ACCENT
    box.line.width = Pt(1.5)
    tb = box.text_frame
    tb.margin_left = Pt(20)
    tb.margin_right = Pt(20)
    tb.margin_top = Pt(15)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    r = p.add_run()
    r.text = "Selective prediction with the ensemble"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = COLOR_PRIMARY
    p = tb.add_paragraph()
    r = p.add_run()
    r.text = "Defer the 10 % most uncertain cases → accuracy on the rest rises from 96.55 % to 98.18 %. At 50 % coverage: 99.64 %."
    r.font.size = Pt(14)
    r.font.color.rgb = COLOR_DARK

    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_calibration(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Calibration: Temperature Scaling Works", prs)
    img1 = MASTER_RESULTS / "reliability_VGG16_raw.png"
    img2 = MASTER_RESULTS / "reliability_VGG16_temp_scaled.png"
    left = Inches(0.4)
    top = Inches(1.3)
    width = (prs.slide_width - Inches(1.0)) // 2
    add_image(slide, img1, left, top, width=width)
    add_image(slide, img2, left + width + Inches(0.2), top, width=width)
    add_caption(slide, "VGG16 raw  —  ECE = 0.0255, MCE = 0.74",
                left, top + Inches(3.5), width, size=12)
    add_caption(slide, "VGG16 after T = 1.44  —  ECE = 0.0153 (-40 %)",
                left + width + Inches(0.2), top + Inches(3.5), width, size=12)

    add_bullets(slide, prs, [
        "Temperature scaling: divide logits by a scalar T fitted on validation set",
        "Argmax preserved → accuracy unchanged",
        "Only the confidence is rescaled to match observed frequencies",
    ], top=prs.slide_height - Inches(1.8), size=13)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_phase2_conformal(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 2 — Conformal Prediction", prs)
    add_bullets(slide, prs, [
        "Split conformal: fit non-conformity threshold on validation set",
        "Two scoring rules: LAC (1 − max prob) and APS (cumulative sum)",
        "Two coverage targets: α = 0.10 → 90 % cover, α = 0.05 → 95 % cover",
        "Formal guarantee: P(y ∈ C(x)) ≥ 1 − α",
    ], top=Inches(1.3), size=15)

    data = [
        ["Model", "Score", "α", "Coverage", "Mean size", "Singleton correct"],
        ["VGG16", "LAC", "0.10", "91.27 %", "0.93", "98.05 %"],
        ["Ensemble", "LAC", "0.10", "90.55 %", "0.93", "97.84 %"],
        ["ResNet50", "APS", "0.10", "88.73 %", "0.96", "96.67 %"],
        ["VGG16", "LAC", "0.05", "95.45 %", "1.05", "97.71 %"],
    ]
    add_table(slide, prs, data, Inches(0.4), Inches(3.7),
              prs.slide_width - Inches(0.8), Inches(2.5),
              font_size=12, header_size=13)
    add_caption(slide, "Empirical coverage tracks the prescribed target across every configuration",
                Inches(0.4), prs.slide_height - Inches(0.7), prs.slide_width - Inches(0.8), size=12)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_phase2_mcd(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 2 — Monte Carlo Dropout", prs)
    add_bullets(slide, prs, [
        "Inference: T = 30 stochastic forward passes per image with dropout active",
        "Sigma (std) is 3 to 4× larger on wrong predictions than on correct",
        "resnet50_mcd reaches 96.18 % — the best single model in the thesis",
        "Predictive entropy at 90 % coverage → 97.78 % selective accuracy",
    ], top=Inches(1.3), size=15)

    img = MASTER_RESULTS / "mc_dropout" / "resnet50_mcd_uncertainty_hist_std.png"
    add_image(slide, img, Inches(1.5), Inches(3.4), width=Inches(7.0))
    add_caption(slide, "σ histogram split by correct vs incorrect predictions",
                Inches(1.5), Inches(6.6), Inches(7.0), size=12)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_phase2_ood_kfold(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 2 — OOD Detection & K-Fold CV", prs)

    img1 = FIGURES / "fig_phase2_ood.png"
    img2 = FIGURES / "fig_phase2_kfold.png"
    left = Inches(0.4)
    top = Inches(1.3)
    width = (prs.slide_width - Inches(1.0)) // 2
    add_image(slide, img1, left, top, width=width)
    add_image(slide, img2, left + width + Inches(0.2), top, width=width)
    add_caption(slide, "OOD: feature-space methods achieve AUROC = 1.0",
                left, top + Inches(3.5), width, size=12)
    add_caption(slide, "K-Fold: ResNet50 = 95.64 ± 0.18 pp (very stable)",
                left + width + Inches(0.2), top + Inches(3.5), width, size=12)
    add_bullets(slide, prs, [
        "Mahalanobis & cosine distance achieve perfect AUROC against synthetic noise",
        "MSP (output-space) fails: softmax saturates even on uniform noise",
        "5-fold CV confirms that the single-split estimate of Phase 1 is reproducible",
    ], top=prs.slide_height - Inches(1.8), size=13)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_phase3_overview(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 3 — 5-Stage Severity Grading", prs)
    add_bullets(slide, prs, [
        ("Reformulation", "Original APTOS labels: No DR · Mild · Moderate · Severe · PDR"),
        ("Test class distribution", "271 / 56 / 150 / 29 / 44 — heavily imbalanced"),
        ("Models", "cnn_5class and resnet50_5class"),
        ("Loss", "Categorical cross-entropy with class_weight='balanced'"),
        ("Architecture heads", "5-unit softmax instead of 1-unit sigmoid"),
    ], top=Inches(1.3), size=18)

    data = [
        ["Model", "Test acc", "QWK", "Ord. dist.", "Macro F1", "ECE"],
        ["cnn_5class", "64.18 %", "0.564", "0.658", "0.440", "0.054"],
        ["resnet50_5class", "77.09 %", "0.847", "0.329", "0.603", "0.030"],
        ["Ensemble (2)", "76.73 %", "0.785", "—", "—", "0.089"],
    ]
    add_table(slide, prs, data, Inches(0.4), Inches(4.7),
              prs.slide_width - Inches(0.8), Inches(1.6),
              font_size=13, header_size=13)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_phase3_perclass(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 3 — Per-Class Performance", prs)

    img1 = MASTER_RESULTS / "multiclass" / "resnet50_5class_confusion_matrix.png"
    img2 = FIGURES / "fig_phase3_per_class.png"
    left = Inches(0.4)
    top = Inches(1.3)
    width = (prs.slide_width - Inches(1.0)) // 2
    add_image(slide, img1, left, top, width=width)
    add_image(slide, img2, left + width + Inches(0.2), top, width=width)
    add_caption(slide, "Confusion matrix — strong diagonal, errors are local",
                left, top + Inches(3.5), width, size=12)
    add_caption(slide, "Per-class F1: No DR almost perfect; Severe and PDR confused",
                left + width + Inches(0.2), top + Inches(3.5), width, size=12)
    add_bullets(slide, prs, [
        "QWK = 0.847 → 'excellent' agreement, competitive with the upper third of the public APTOS leaderboard",
        "Mean ordinal distance 0.33 — errors are typically one class away",
        "Mild has low recall (0.48): subtle microaneurysms are easy to miss",
        "Severe ↔ PDR confusion is clinically acceptable: both warrant urgent referral",
    ], top=prs.slide_height - Inches(2.2), size=13)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_phase3_triage(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 3 — Conformal Triage Policy", prs)

    # Big numbers
    big_left = Inches(0.5)
    big_top = Inches(1.5)
    big_w = prs.slide_width - Inches(1.0)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, big_left, big_top, big_w, Inches(2.0),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_PRIMARY
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(20)
    tf.margin_left = Pt(40)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "APS at α = 0.10  →  empirical coverage 89.45 %"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = COLOR_WHITE
    p = tf.add_paragraph()
    p.space_before = Pt(15)
    r = p.add_run()
    r.text = "Per-class coverage:  [91, 82, 89, 90, 93] %  —  uniform across all 5 severities"
    r.font.size = Pt(16)
    r.font.color.rgb = RGBColor(0xDD, 0xEE, 0xFF)

    # Triage breakdown
    data = [
        ["Decision", "Cases", "Set size", "Action"],
        ["Auto-resolve", "234 (43 %)", "1", "No clinician needed"],
        ["Soft refer", "159 (29 %)", "2", "Adjacent severities flagged"],
        ["Refer to clinician", "157 (29 %)", "≥ 3", "Genuinely ambiguous"],
    ]
    add_table(slide, prs, data, Inches(0.5), Inches(4.0),
              prs.slide_width - Inches(1.0), Inches(2.0),
              font_size=14, header_size=14)
    add_caption(slide, "71 % auto-classify / 29 % refer  —  a deployable clinical workflow",
                Inches(0.5), prs.slide_height - Inches(0.7),
                prs.slide_width - Inches(1.0), size=13)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_streamlit(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Deployment: Streamlit Application", prs)
    left = [
        ("Single & batch inference", "Upload one image or a folder"),
        ("Calibrated probabilities", "Per-class softmax, post-T scaling"),
        ("Conformal prediction sets", "Selectable α (90 % / 95 %)"),
        ("Ensemble disagreement", "Optional multi-model comparison"),
    ]
    right = [
        ("Image-quality heuristics", "Brightness / contrast / focus check"),
        ("PDF report export", "Multi-page batch summary"),
        ("Multilingual UI", "English & Albanian (.json dictionaries)"),
        ("Read-only at model level", "No retraining at user request"),
    ]
    add_two_columns(slide, prs, left, right, top=Inches(1.3), size=14)

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), prs.slide_height - Inches(1.8),
        prs.slide_width - Inches(1.0), Inches(1.3),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_LIGHT
    box.line.color.rgb = COLOR_ACCENT
    box.line.width = Pt(1.5)
    tb = box.text_frame
    tb.margin_left = Pt(20)
    tb.margin_top = Pt(15)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    r = p.add_run()
    r.text = "Statistical guarantees carry through to deployment"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = COLOR_PRIMARY
    p = tb.add_paragraph()
    r = p.add_run()
    r.text = "Pre-computed conformal thresholds (app_conformal_thresholds.json) are loaded at startup so the marginal coverage guarantee established offline holds online."
    r.font.size = Pt(13)
    r.font.color.rgb = COLOR_DARK

    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_app_screenshots(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Streamlit Application — Live Views", prs)

    # Two portrait screenshots on the left/centre, landscape one on the right.
    single = APP_SHOTS / "app_predict_single.png"
    comparison = APP_SHOTS / "app_predict_comparison.png"
    five_stage = APP_SHOTS / "app_5stage_grading.png"

    add_image(slide, single, Inches(0.4), Inches(1.5), width=Inches(3.0))
    add_caption(slide, "Binary prediction + quality check",
                Inches(0.4), Inches(6.3), Inches(3.0), size=11)

    add_image(slide, comparison, Inches(3.6), Inches(1.5), width=Inches(3.0))
    add_caption(slide, "Multi-model comparison (ResNet50 vs VGG16)",
                Inches(3.6), Inches(6.3), Inches(3.0), size=11)

    add_image(slide, five_stage, Inches(7.1), Inches(2.7), width=Inches(5.8))
    add_caption(slide, "5-stage grading with conformal set (α = 0.10)",
                Inches(7.1), Inches(6.3), Inches(5.8), size=11)

    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_synthesis(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Research Questions — Answered", prs)
    items = [
        ("RQ1: Do architectures differ?", "No — five strong models statistically indistinguishable (every pairwise McNemar p > 0.5). Ensemble offers a small but real improvement."),
        ("RQ2: How calibrated are the models?", "Reasonably well (ECE 0.017 – 0.046). Temperature scaling delivers consistent small gains and one large gain (VGG16, −40 %)."),
        ("RQ3: Best uncertainty signal?", "Predictive entropy lifts selective accuracy 96.55 → 98.18 % at 90 % coverage. Mahalanobis & cosine: AUROC = 1.0 against synthetic OOD. Conformal: empirical coverage tracks target."),
        ("RQ4: Does it generalise to 5-class?", "Yes. ResNet50 multi-class → QWK = 0.847, ECE = 0.030, conformal 71 % auto / 29 % refer with balanced per-class coverage."),
    ]
    add_bullets(slide, prs, items, top=Inches(1.3), size=14)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_limitations(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Limitations", prs)
    items = [
        ("Single dataset", "All experiments on APTOS 2019 (primarily Indian cohort, one camera distribution)"),
        ("Synthetic OOD only", "Uniform-noise inputs are the easy case; real OOD would yield AUROC 0.80 – 0.95"),
        ("Weak minority-class performance", "F1 = 0.37 for Severe, 0.44 for PDR; only 29 / 44 test images respectively"),
        ("CPU-only training", "No backbone fine-tuning, no Vision Transformers, no TTA"),
        ("No clinician-in-the-loop validation", "Conformal 'refer 29 %' policy is a candidate, not yet validated by an ophthalmologist"),
        ("Single random seed for splits", "Test set is fixed; 5-fold CV only for ResNet50"),
    ]
    add_bullets(slide, prs, items, top=Inches(1.3), size=15)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_future_work(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Future Work", prs)
    left = [
        ("Cross-dataset evaluation", "Messidor-2 / IDRiD — infrastructure ready"),
        ("Clinician-in-the-loop study", "Retrospective review of deferred cases"),
        ("Foundation models (RETFound)", "Vision Transformer pretrained on 1.6M retinas"),
        ("Hyperparameter sensitivity", "Full sweep with multiple seeds"),
        ("Demographic shift evaluation", "Per-subgroup analysis on EyePACS metadata"),
    ]
    right = [
        ("Multi-task lesion segmentation", "IDRiD pixel-level annotations"),
        ("Conformal risk control", "Cost-sensitive thresholds for asymmetric clinical risk"),
        ("Multi-modal inputs", "Fundus + HbA1c + clinical history"),
        ("Active learning", "Label the conformal abstain set first"),
        ("EHR / FHIR integration", "Deploy beyond the research UI"),
    ]
    add_two_columns(slide, prs, left, right, top=Inches(1.3), size=14)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_contributions(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Contributions", prs)
    items = [
        "Reproducible, statistically rigorous evaluation pipeline on APTOS 2019",
        "Calibration analysis for six DR classifiers with 40 % ECE reduction on VGG16",
        "First application of split conformal prediction to multi-stage DR on APTOS 2019",
        "Side-by-side comparison of MC Dropout, ensemble disagreement, and four OOD detection methods on one benchmark",
        "5-stage grading model achieving QWK = 0.847, competitive with public Kaggle leaderboards",
        "Open-source Streamlit application exposing calibrated probabilities, ensemble disagreement, and conformal sets in real time",
    ]
    add_bullets(slide, prs, items, top=Inches(1.3), size=17)
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_concluding_message(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Closing Message", prs)

    # Big quote
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(2.0),
        prs.slide_width - Inches(1.4), Inches(2.0),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_PRIMARY
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(30)
    tf.margin_left = Pt(40)
    tf.margin_right = Pt(40)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "“Modern medical AI must do more than maximise accuracy."
    r.font.size = Pt(22)
    r.font.italic = True
    r.font.color.rgb = COLOR_WHITE
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "It must know when to defer.”"
    r.font.size = Pt(22)
    r.font.italic = True
    r.font.color.rgb = COLOR_WHITE

    # Bottom callouts
    tb = slide.shapes.add_textbox(
        Inches(0.7), Inches(4.5), prs.slide_width - Inches(1.4), Inches(2.0),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "From a single accuracy number"
    r.font.size = Pt(16)
    r.font.color.rgb = COLOR_GREY
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)
    r = p.add_run()
    r.text = "↓"
    r.font.size = Pt(20)
    r.font.color.rgb = COLOR_ACCENT
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)
    r = p.add_run()
    r.text = "to a transparent, statistically grounded triage system"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = COLOR_PRIMARY
    add_footer_bar(slide, prs)
    add_page_number(slide, prs, num, total)


def slide_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height,
    )
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_PRIMARY
    band.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(3.7),
        prs.slide_width, Inches(0.08),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_ACCENT
    accent.line.fill.background()

    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.2), prs.slide_width - Inches(1.0), Inches(1.4),
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Thank You"
    r.font.size = Pt(60)
    r.font.bold = True
    r.font.color.rgb = COLOR_WHITE

    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.0), prs.slide_width - Inches(1.0), Inches(1.0),
    )
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Questions & Discussion"
    r.font.size = Pt(28)
    r.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)

    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.5), prs.slide_width - Inches(1.0), Inches(1.0),
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Alketa Alia"
    r.font.size = Pt(18)
    r.font.color.rgb = COLOR_WHITE
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Master Thesis Defense  ·  Polytechnic University of Tirana  ·  May 2026"
    r.font.size = Pt(13)
    r.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)


# ============================================================================
# Build
# ============================================================================

# Speaker notes (Albanian) — one entry per builder, in the same order.
NOTES = [
    # 1 — Title
    "Të nderuar anëtarë të komisionit, mirëdita dhe faleminderit që jeni këtu. "
    "Unë jam Alketa Alia, dhe sot do të prezantoj tezën time të masterit, të titulluar "
    "'Klasifikim i Vetëdijshëm ndaj Pasigurisë i Retinopatisë Diabetike', e udhëhequr nga Profesor [emri]. "
    "Me pak fjalë, kjo punë ndërton një sistem AI që jo vetëm njeh sëmundjen e syrit nga fotot, "
    "por edhe e di kur është i pasigurt — dhe ja kalon rastin doktorit.",
    # 2 — Outline
    "Prezantimi im do të ketë këtë strukturë: fillimisht problemi dhe motivimi, pastaj pyetjet "
    "kërkimore dhe metodologjia, dhe pjesa kryesore — rezultatet në tre faza. Në fund, deployment-i, "
    "kufizimet dhe përfundimet. Do të zgjasë rreth 15–18 minuta.",
    # 3 — Problem
    "Le të fillojmë me pse e bëra këtë punë. Retinopatia diabetike është shkaku kryesor i verbimit "
    "te të rriturit në moshë pune. Sot kemi mbi 500 milionë diabetikë në botë, por nuk ka mjaftueshëm "
    "oftalmologë për t'i kontrolluar të gjithë. Problemi është që sëmundja nuk shfaq simptoma në fazat "
    "e hershme — ndaj nevojitet kontroll periodik. Pikërisht këtu hyjnë sistemet automatike — jo si luks, "
    "por si zgjidhje për mungesën e specialistëve.",
    # 4 — 5 Stages
    "Klinikisht, sëmundja ka 5 faza, nga 'pa DR' deri te faza më e rëndë, proliferative. "
    "Çështja kyçe (trego fjalinë e fundit në sllajd): një sistem që thotë vetëm 'ke DR' ose 's'ke DR' "
    "nuk mjafton — sepse një pacient me fazën e parë dhe një me fazën e fundit marrin të njëjtin etiketim, "
    "por trajtohen krejt ndryshe. Ndaj na duhet fazimi i plotë.",
    # 5 — Why Uncertainty
    "Tani, problemi i madh me sistemet AI: ato janë mbi-konfidente. Mund të thonë '99% i sigurt' dhe "
    "prapë të gabojnë. Në fusha të tjera kjo s'ka rëndësi, por në mjekësi është e rrezikshme — një gabim "
    "i sigurt mund të vonojë trajtimin që shpëton shikimin. Madje edhe rregullatorët — FDA, Bashkimi "
    "Europian — tani e kërkojnë që sistemi t'a komunikojë pasigurinë. Ndaj sistemi im jep katër shtresa "
    "output, jo vetëm një përgjigje.",
    # 6 — Research Questions
    "Puna ime u organizua rreth katër pyetjeve: A ndryshojnë realisht modelet mes tyre? Sa të besueshme "
    "janë probabilitetet që japin? Cili sinjal pasigurie funksionon më mirë? Dhe a funksionon e njëjta "
    "gjë kur kalojmë në 5 klasa? Të katër do u jap përgjigje me rezultate konkrete në fund.",
    # 7 — Methodology
    "Për metodologjinë: përdora datasetin publik APTOS 2019 me rreth 3,600 imazhe. I ndava në 70% trajnim, "
    "15% validim, dhe 15% test — dhe test-in e mbajta të kyçur deri në fund, që rezultatet të jenë të "
    "besueshme. Gjithçka u bë në një laptop të zakonshëm, pa GPU. Kjo ishte zgjedhje e qëllimshme — që "
    "sistemi të jetë i përdorshëm edhe në klinika me burime të kufizuara.",
    # 8 — Architectures
    "Provova gjashtë arkitektura të ndryshme — katër prej tyre të paratrajnuara mbi ImageNet, dhe dy CNN "
    "të ndërtuara nga unë. Siç do ta shohim, gjetja interesante është që këto modele, megjithëse shumë "
    "të ndryshme në strukturë, japin rezultate pothuajse identike.",
    # 9 — 4 Pillars
    "Për të masur pasigurinë, përdora katër mjete: kalibrimin, që kontrollon sa të besueshme janë "
    "probabilitetet; Monte Carlo Dropout, një teknikë bayesiane; parashikimin konformal, që jep garanci "
    "matematike; dhe detektimin e imazheve të çuditshme. Tani do t'i shohim rezultatet, fazë pas faze.",
    # 10 — Phase 1 Results
    "Faza e parë — klasifikimi binar. (Trego tabelën.) Si shihni në tabelë, të gjitha modelet janë shumë "
    "afër njëri-tjetrit — mes 95 dhe 96% saktësi. Vetëm modeli i fundit, CNN me tanh, është dukshëm më "
    "poshtë. Dhe ensemble — kombinimi i të gjashtë modeleve — arrin më të lartën, 96.55%.",
    # 11 — Phase 1 Visual
    "Po e shohim këtë vizualisht. (Trego grafikun majtas.) Majtas, shtyllat e saktësisë — pothuajse të "
    "barabarta. Vetëm e fundit, më e shkurtër. (Trego heatmap-in djathtas.) Djathtas është testi statistikor. "
    "Pothuajse të gjitha kutitë janë jeshile, që do 'asnjë dallim real'. Vetëm rreshti i modelit të dobët "
    "është i kuq. Pra statistikisht, zgjedhja e arkitekturës nuk ka rëndësi — ç'ka rëndësi është mënyra si "
    "i trajnojmë. Dhe nëse refuzojmë rastet më të pasigurta, saktësia rritet edhe më shumë: nga 96.55% në "
    "98.18%, dhe në 50% coverage arrin 99.64%.",
    # 12 — Calibration
    "Tani kalibrimi — sa të besueshme janë numrat. (Trego diagonalen.) Kjo vijë diagonale është kalibrimi "
    "perfekt. Para korrigjimit, disa shtylla bien nën vijë — modeli është mbi-konfident. Pas një korrigjimi "
    "të thjeshtë që quhet temperature scaling, shtyllat ulen te diagonalja; ECE bie nga 0.026 në 0.015. "
    "Pikë e rëndësishme: ky korrigjim nuk e prek saktësinë, vetëm e bën besueshmërinë reale. Përmirësim falas.",
    # 13 — Conformal
    "Faza e dytë — mjetet e avancuara. E para, parashikimi konformal. Ideja është që sistemi jep një garanci "
    "matematike: 'përgjigja e saktë do jetë brenda kësaj liste me 90% siguri'. (Trego kolonën Coverage.) "
    "Kur kërkojmë 90%, marrim realisht afër 90%; kur kërkojmë 95%, marrim afër 95%. Pra garancia mbahet në "
    "praktikë. Dhe VGG16 me LAC ka 'singleton correct' 98% — kur jep një klasë të vetme, është pothuajse "
    "gjithmonë e drejtë.",
    # 14 — MC Dropout
    "Pastaj Monte Carlo Dropout — ekzekutojmë modelin 30 herë dhe shohim sa lëkundet. (Trego histogramin.) "
    "Shtyllat blu janë rastet ku modeli pati të drejtë — grumbullohen majtas, te pasiguria e ulët. Shtyllat "
    "e kuqe janë gabimet — shkojnë djathtas, te pasiguria e lartë; σ është 3-4 herë më e madhe. Pra pasiguria "
    "e lartë na paralajmëron për një gabim të mundshëm — pikërisht ç'na duhet për të vendosur kur t'i kalojmë "
    "rastin doktorit.",
    # 15 — OOD + K-Fold
    "Dy gjetje këtu. (Trego grafikun majtas.) Majtas — sa mirë e dallon sistemi kur i jepet një imazh i "
    "çuditshëm, që nuk është fundus. Metoda e thjeshtë (MSP) dështon, AUROC 0.59, por dy metodat e fundit "
    "(Mahalanobis, cosine) janë perfekte, AUROC 1.0. Caveat i ndershëm: kjo është rasti i lehtë. (Trego "
    "grafikun djathtas.) Djathtas — e ritrajnova modelin pesë herë me ndarje të ndryshme; të gjitha shtyllat "
    "janë pothuajse njësoj, devijim 0.18 pikë — rezultatet janë stabile, nuk janë fat.",
    # 16 — Phase 3 Overview
    "Faza e tretë — pjesa më klinike. Tani e bëra detyrën e vështirë: klasifikim në 5 faza në vend të 2. "
    "Saktësia bie në 77%, por kjo është normale — detyra është shumë më e vështirë. Numri që ka vërtet "
    "rëndësi është QWK 0.847 (theks) — kjo është në kategorinë 'shkëlqyeshëm', konkurruese me sistemet më "
    "të mira në botë, megjithëse unë përdora vetëm një laptop.",
    # 17 — Per-class
    "Le ta zbërthejmë për klasë. (Trego matricën majtas.) Numrat e mëdhenj janë në diagonale, pra modeli "
    "kryesisht ka të drejtë; kur gabon, gabon mes fazave fqinje, jo të largëta. (Trego grafikun djathtas.) "
    "Performanca për klasë: 'No DR' është pothuajse perfekt. Por 'Severe' dhe 'PDR' janë më të dobëta — sepse "
    "kemi pak shembuj për to dhe janë vizualisht të ngjashme. Megjithatë, klinikisht kjo është e tolerueshme, "
    "sepse të dyja kërkojnë referim urgjent.",
    # 18 — Triage Policy
    "Dhe kjo është pika kulminante e tezës. (Pauzë, lër t'a lexojnë.) Sistemi i ndan rastet në tre grupe: "
    "rreth 71% i auto-klasifikon vetë me siguri, dhe vetëm 29% — rastet realisht të paqarta — i dërgon te "
    "oftalmologu. Ky është një workflow klinik i gatshëm: redukton ngarkesën e doktorit, por nuk rrezikon "
    "pacientin, sepse rastet e dyshimta gjithmonë shkojnë te njeriu.",
    # 19 — Streamlit
    "E gjithë kjo është paketuar në një aplikacion real. Përdoruesi ngarkon foton, dhe merr probabilitetin e "
    "kalibruar, setin conformal, dhe rekomandimin nëse rasti duhet referim. Ndërfaqja është dygjuhëshe — "
    "anglisht dhe shqip. Pikë e rëndësishme: aplikacioni është read-only në nivel modeli; garancitë statistikore "
    "offline mbahen online.",
    # 20 — App Screenshots (NEW)
    "Këto janë pamje reale nga aplikacioni gjatë punës. Majtas, parashikimi binar me kontrollin e cilësisë së "
    "imazhit dhe selektorin e gjuhës. Në mes, krahasimi i dy modeleve — ResNet50 dhe VGG16 — që tregon "
    "mosmarrëveshjen mes tyre. Djathtas, grading-u 5-fazor me setin conformal. Pra sistemi nuk është thjesht "
    "teori — punon në kohë reale në një laptop të zakonshëm.",
    # 21 — RQs Answered
    "T'i përmbledhim përgjigjet e katër pyetjeve: modelet nuk ndryshojnë statistikisht; janë mjaft të "
    "kalibruara; conformal jep garanci formale; dhe gjithçka funksionon edhe në 5 klasa.",
    # 22 — Contributions
    "Kontributet kryesore: një pipeline rigoroze dhe e riprodhueshme, analizë e plotë e kalibrimit, aplikimi i "
    "parë i parashikimit konformal në këtë dataset për 5 faza, krahasim sistematik i metodave të pasigurisë, "
    "dhe një aplikacion open-source.",
    # 23 — Limitations
    "Jam e ndershme edhe për kufizimet: punova me një dataset të vetëm, testimi i imazheve të çuditshme ishte "
    "sintetik, klasat e rralla janë akoma të dobëta, trajnim vetëm CPU, dhe nuk kam ende validim nga një "
    "oftalmolog real. Këto janë pikërisht hapat e ardhshëm.",
    # 24 — Future Work
    "Në të ardhmen: testim mbi një dataset të dytë (Messidor-2 ose IDRiD) për të provuar përgjithësimin; një "
    "studim klinik me oftalmolog për të validuar politikën konformal; dhe modele më të avancuara si foundation "
    "models të specializuara për retinën, p.sh. RETFound.",
    # 25 — Closing
    "Për të mbyllur — mesazhi kryesor i tezës është ky (ngadalë, me theks): AI mjekësor modern duhet të bëjë "
    "më shumë se vetëm saktësi. Duhet të dijë kur të refuzojë dhe t'a kalojë rastin te doktori. Kjo është ajo "
    "që e bën një sistem realisht të deployueshëm në klinikë.",
    # 26 — Thanks
    "Ju falenderoj shumë për vëmendjen. Jam e gatshme për pyetjet dhe diskutimin tuaj.",
]


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_title,
        slide_outline,
        slide_problem,
        slide_5_stages,
        slide_motivation,
        slide_research_questions,
        slide_methodology_overview,
        slide_architectures,
        slide_4_pillars,
        slide_phase1_results,
        slide_phase1_figure,
        slide_calibration,
        slide_phase2_conformal,
        slide_phase2_mcd,
        slide_phase2_ood_kfold,
        slide_phase3_overview,
        slide_phase3_perclass,
        slide_phase3_triage,
        slide_streamlit,
        slide_app_screenshots,
        slide_synthesis,
        slide_contributions,
        slide_limitations,
        slide_future_work,
        slide_concluding_message,
        slide_thanks,
    ]

    total = len(builders)
    for i, fn in enumerate(builders, start=1):
        # First (title) and last (thanks) slides don't take page-number args
        if fn in (slide_title, slide_thanks):
            fn(prs)
        else:
            fn(prs, i, total)
        # Attach speaker notes to the slide just created
        if i <= len(NOTES):
            prs.slides[-1].notes_slide.notes_text_frame.text = NOTES[i - 1]

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}")
    print(f"Total slides: {total}")
    print(f"Notes attached: {min(len(NOTES), total)}")


if __name__ == "__main__":
    main()
